"""
Build a synthetic .pptx that represents the "deck-with-implicit-pattern"
case — every slide on the default layout, content as hand-positioned text
boxes, visual identity in the SLIDES not in any master/layout.

Use as a test fixture for the classifier (should report
"deck-with-implicit-pattern") and for the synthesizer (should cluster the
recurring visual patterns into derived layouts).

Output: tests/fixtures/synthetic-deck-no-template.pptx

The fixture contains 4 visual patterns repeated:
    Pattern A — title slide (1 big text top-left)
    Pattern B — three-column layout (3 text boxes side-by-side)
    Pattern C — big number ("metrics" feel)
    Pattern D — quote (large centered text + attribution)

Each pattern repeats 2x so the cluster step has signal to work with.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


# Visual identity — consistent across all slides, but applied per-shape
# (not via master/layout). The synthesizer should recover these as tokens.
BRAND_DARK = RGBColor(0x1F, 0x2A, 0x44)
BRAND_LIGHT = RGBColor(0x4A, 0x51, 0x60)
ACCENT = RGBColor(0xE2, 0xA0, 0x3F)
FONT_FAMILY = "Helvetica"


def add_text_box(slide, left_in, top_in, width_in, height_in, text, *,
                  size_pt=18, bold=False, color=BRAND_DARK):
    """Add a hand-positioned text box. This is what makes the slide a
    'deck' rather than using a layout placeholder."""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_pattern_a_title(prs, title_text, subtitle_text):
    """Pattern A — title slide. 1 large title, 1 small subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = Blank
    add_text_box(slide, 0.7, 2.3, 9.0, 1.5, title_text,
                 size_pt=48, bold=True, color=BRAND_DARK)
    add_text_box(slide, 0.7, 4.0, 9.0, 0.6, subtitle_text,
                 size_pt=18, color=BRAND_LIGHT)


def add_pattern_b_three_column(prs, header, col1, col2, col3):
    """Pattern B — three-column layout. 1 header + 3 column text boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_text_box(slide, 0.7, 0.6, 9.0, 0.8, header,
                 size_pt=28, bold=True, color=BRAND_DARK)
    col_width = 2.7
    col_top = 1.8
    col_height = 4.5
    for i, col_text in enumerate([col1, col2, col3]):
        left = 0.7 + i * (col_width + 0.2)
        add_text_box(slide, left, col_top, col_width, col_height, col_text,
                     size_pt=14, color=BRAND_DARK)


def add_pattern_c_metric(prs, value_text, label_text):
    """Pattern C — big number. 1 huge value, 1 small label below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_text_box(slide, 0.7, 1.8, 9.0, 2.5, value_text,
                 size_pt=120, bold=True, color=ACCENT)
    add_text_box(slide, 0.7, 4.5, 9.0, 0.8, label_text,
                 size_pt=20, color=BRAND_DARK)


def add_pattern_d_quote(prs, quote_text, attribution_text):
    """Pattern D — quote slide. Large centered text + attribution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_text_box(slide, 1.0, 2.0, 8.0, 2.5, quote_text,
                 size_pt=28, bold=False, color=BRAND_DARK)
    add_text_box(slide, 1.0, 4.7, 8.0, 0.6, attribution_text,
                 size_pt=14, color=BRAND_LIGHT)


def build(out_path: Path) -> None:
    prs = Presentation()  # default 16:9, all default layouts (Title Slide, etc.)

    # Pattern A x2
    add_pattern_a_title(prs, "Q1 portfolio review",
                        "Where we are, where we're going.")
    add_pattern_a_title(prs, "Q2 outlook",
                        "Three commitments and a wildcard.")

    # Pattern B x2
    add_pattern_b_three_column(prs,
        "What's working",
        "Pipeline velocity\nUp 18% QoQ. Sales cycle compressed by 9 days.",
        "Customer NPS\nHolding at 58. Inbound expansion outpacing churn.",
        "Hiring\nClosed 7 of 9 open roles. Strong inbound.",
    )
    add_pattern_b_three_column(prs,
        "Three risks",
        "Concentration\nTop 5 customers = 41% of revenue. Single dependency.",
        "Org load\nEng + design carrying 2x prior load. Burnout risk.",
        "Competitive\nTwo well-funded entrants. 12-month window.",
    )

    # Pattern C x2
    add_pattern_c_metric(prs, "47%", "growth in active accounts, year over year")
    add_pattern_c_metric(prs, "$8.2M", "in net new ARR this quarter")

    # Pattern D x2
    add_pattern_d_quote(prs,
        "I knew the date picker existed. I just couldn't tell from the docs whether it handled our case. So I rebuilt it.",
        "Senior engineer, consumer team",
    )
    add_pattern_d_quote(prs,
        "The clearest signal from the user research was that no one read the onboarding email.",
        "Product manager, growth",
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path} with {len(prs.slides)} slides across 4 visual patterns.")


def main() -> int:
    parser = argparse.ArgumentParser(description="Build the deck-with-implicit-pattern test fixture.")
    parser.add_argument(
        "out", nargs="?",
        default="tests/fixtures/synthetic-deck-no-template.pptx",
        help="Output .pptx path",
    )
    args = parser.parse_args()
    build(Path(args.out).expanduser())
    return 0


if __name__ == "__main__":
    sys.exit(main())
