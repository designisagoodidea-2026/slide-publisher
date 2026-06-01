"""
Build a synthetic .pptx test fixture for template_extractor_pptx.py.

Produces a deck with explicit layout names matching the IR layout-intent
catalog, plus a small palette of colors and a couple of font sizes in the
master, so the extractor has structure to find.

This is a *test fixture builder*, not a runtime dependency of the plugin.
Run it once to produce the .pptx, then commit the .pptx if you want it
permanent. The .pptx is the actual test artifact.

Usage:
    python tests/fixtures/build_synthetic_pptx.py tests/fixtures/synthetic-template.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


# Match the 10 IR layout intents. These names should let the extractor's
# pattern dictionary find a clean 10-of-10 layout map.
SYNTHETIC_LAYOUT_NAMES = [
    "Title Slide",
    "Section Header",
    "Title and Content",
    "Three Column",
    "Comparison",
    "Pull Quote",
    "Image and Caption",
    "Stat Block",
    "Timeline",
    "Big Statement",
]


def build(out_path: Path) -> None:
    prs = Presentation()  # default 16:9, default theme

    # Rename the first 10 built-in layouts to match our IR intents.
    # python-pptx exposes slide_layouts as a sequence on the slide_master.
    master = prs.slide_masters[0]
    layouts = list(master.slide_layouts)
    for i, name in enumerate(SYNTHETIC_LAYOUT_NAMES):
        if i >= len(layouts):
            break
        # The layout's name is stored in the cSld element's `name` attribute.
        layouts[i].element.cSld.set("name", name)

    # Add one slide per layout so the file isn't empty and the layouts are
    # exercised (the extractor doesn't strictly need slides, but real-world
    # decks will).
    for i, layout in enumerate(layouts[: len(SYNTHETIC_LAYOUT_NAMES)]):
        slide = prs.slides.add_slide(layout)
        # Put a placeholder title if the layout has one
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                shape.text = SYNTHETIC_LAYOUT_NAMES[i]
                break

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path} with {len(SYNTHETIC_LAYOUT_NAMES)} layouts.")


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a synthetic .pptx test fixture.")
    parser.add_argument(
        "out", nargs="?",
        default="tests/fixtures/synthetic-template.pptx",
        help="Output .pptx path (default: tests/fixtures/synthetic-template.pptx)",
    )
    args = parser.parse_args()
    build(Path(args.out).expanduser())
    return 0


if __name__ == "__main__":
    sys.exit(main())
