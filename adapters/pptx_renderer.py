"""
pptx_renderer.py — render a validated Deck IR into a .pptx using the user's
template profile.

Inputs:
    ir_path:       path to a Deck IR YAML conforming to ir/schema.json
    profile_path:  path to a template profile YAML conforming to
                   template-profile/schema.json (must have a `templates.pptx`
                   branch)
    out_path:      where to write the rendered .pptx

Outputs:
    <out_path>            — the rendered .pptx
    <out_path>.loss.md    — markdown loss manifest
    <out_path>.loss.json  — JSON sidecar (machine-readable)

Loss-manifest categories (mirror translation-engine convention):
    LOSSLESS   — preserved exactly
    LOSSY      — preserved with degradation (named)
    DROPPED    — not preserved; reason captured
    ANNOTATED  — added by renderer (not in source)

Anonymity: ships in the public plugin. No hard-coded layouts, no user data
defaults. The renderer reads only the user's IR + profile + template.
"""

from __future__ import annotations

import argparse
import datetime as dt
import json
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

try:
    import yaml
except ImportError:
    print("ERROR: pyyaml is not installed. Install with: pip install pyyaml",
          file=sys.stderr)
    sys.exit(2)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx is not installed. Install with: "
          "pip install python-pptx", file=sys.stderr)
    sys.exit(2)


# ---------------------------------------------------------------------------
# Loss manifest — uses _common.LossManifest with format-specific `extra`
# ---------------------------------------------------------------------------

sys.path.insert(0, str(Path(__file__).parent))
from _common import LossManifest, LossEntry, LAYOUT_FALLBACKS as _COMMON_FALLBACKS, check_title_overflow  # noqa: E402

# Stub out the legacy class block so the remaining adapter code (which uses
# LossManifest with template_path) keeps working. We map template_path → extra.
_LEGACY_KEEP_FOR_CODE_BELOW = True
if False:
    @dataclass
    class _LossEntry:  # unreachable; preserved for reader navigation
        category: str
        slide_id: str | None
        field: str
        detail: str

        def to_dict(self) -> dict:
            return {
                "category": self.category,
                "slide_id": self.slide_id,
                "field": self.field,
                "detail": self.detail,
            }


# Shim: existing code constructs LossManifest with `template_path=`; route
# that into the shared dataclass's `extra` dict.
def _make_manifest(*, deck_title: str, rendered_at: str,
                   template_path: str = "") -> LossManifest:
    return LossManifest(
        deck_title=deck_title,
        rendered_at=rendered_at,
        renderer="render-pptx",
        extra={"template_path": template_path},
    )


# Provide attribute access shim so legacy `manifest.template_path` reads work
def _patch_loss_manifest_template_path() -> None:
    def _get(self):  # type: ignore
        return self.extra.get("template_path", "")
    def _set(self, value):  # type: ignore
        self.extra["template_path"] = value
    if not hasattr(LossManifest, "template_path"):
        LossManifest.template_path = property(_get, _set)  # type: ignore


_patch_loss_manifest_template_path()


def _legacy_loss_to_json_with_template(self) -> dict:  # noqa
    """Preserved here as documentation of the legacy shape; not used."""
    return {
            "deck_title": self.deck_title,
            "rendered_at": self.rendered_at,
            "renderer": self.renderer,
            "template_path": self.template_path,
            "summary": self.counts(),
            "entries": [e.to_dict() for e in self.entries],
        }


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def load_ir(path: Path) -> dict:
    return yaml.safe_load(path.read_text())


def load_profile(path: Path) -> dict:
    return yaml.safe_load(path.read_text())


def find_layout(prs: Presentation, layout_name: str) -> Any | None:
    """Find a slide layout by name across all masters. None if not found."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout
    return None


# Fallback chain — if the user's profile doesn't have a mapping for an
# intent, try these alternatives. First found wins. If none of the
# alternatives are in the profile either, fall back to "claim_with_evidence"
# (the universal default), and if THAT's missing, use the first available
# layout in the template.
LAYOUT_FALLBACKS: dict[str, list[str]] = {
    "title": ["callout", "claim_with_evidence"],
    "section_break": ["callout", "claim_with_evidence"],
    "claim_with_evidence": ["callout"],
    "three_pillars": ["claim_with_evidence", "callout"],
    "comparison": ["claim_with_evidence", "three_pillars"],
    "quote": ["callout", "claim_with_evidence"],
    "image_with_caption": ["claim_with_evidence"],
    "metrics": ["three_pillars", "claim_with_evidence"],
    "timeline": ["claim_with_evidence", "three_pillars"],
    "callout": ["claim_with_evidence"],
}


def resolve_layout(
    prs: Presentation,
    layout_map: dict[str, str],
    intent: str,
    manifest: LossManifest,
    slide_id: str,
) -> Any:
    """Resolve the actual pptx layout for a given IR layout intent.

    Walks the fallback chain. Records every substitution as a LOSSY entry
    in the manifest.
    """
    # 1. Direct mapping
    layout_name = layout_map.get(intent)
    if layout_name:
        layout = find_layout(prs, layout_name)
        if layout:
            return layout
        manifest.add(
            "LOSSY", slide_id, "layout_intent",
            f"profile mapped '{intent}' to layout '{layout_name}', "
            f"but that layout was not found in the template.",
        )

    # 2. Fallback chain
    for alt_intent in LAYOUT_FALLBACKS.get(intent, []):
        alt_name = layout_map.get(alt_intent)
        if not alt_name:
            continue
        layout = find_layout(prs, alt_name)
        if layout:
            manifest.add(
                "LOSSY", slide_id, "layout_intent",
                f"intent '{intent}' substituted to '{alt_intent}' "
                f"(template layout '{alt_name}') via the documented "
                f"fallback chain.",
            )
            return layout

    # 3. Universal default
    default_name = layout_map.get("claim_with_evidence")
    if default_name:
        layout = find_layout(prs, default_name)
        if layout:
            manifest.add(
                "LOSSY", slide_id, "layout_intent",
                f"intent '{intent}' had no mapping or fallback; defaulted to "
                f"'claim_with_evidence' (layout '{default_name}').",
            )
            return layout

    # 4. First available layout
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            manifest.add(
                "LOSSY", slide_id, "layout_intent",
                f"intent '{intent}' had no mapping; defaulted to first "
                f"available layout '{layout.name}'.",
            )
            return layout

    raise RuntimeError(f"No layouts available in the template at all.")


def set_title(
    slide: Any, title: str, manifest: LossManifest, slide_id: str,
    intent: str = "claim_with_evidence",
) -> None:
    """Set the slide's title placeholder. If none, add an ANNOTATED entry.

    Also runs a title-overflow check (v0.2.2): when the IR title exceeds the
    layout intent's recommended character budget, the manifest gets an
    ANNOTATED entry so users (and the visual-QA anomaly classifier) can
    pre-empt cramped/wrapped/colliding rendering.
    """
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        check_title_overflow(title, intent, manifest, slide_id)
        return
    manifest.add(
        "DROPPED", slide_id, "title",
        f"layout has no title placeholder; IR title '{title}' was dropped. "
        f"Consider mapping this layout to a different intent or adding a "
        f"title placeholder to the layout in the template.",
    )


def populate_body(
    slide: Any,
    body_blocks: list[dict],
    manifest: LossManifest,
    slide_id: str,
) -> None:
    """Populate a slide's body content from IR body_blocks.

    Strategy (v0.2.1 — multi-placeholder aware):

      1. Classify non-title placeholders by capability: text-capable vs picture-only.
      2. If the layout has zero text placeholders → DROP all body_blocks (record).
      3. Special case — single bullets block + multiple text placeholders:
         split the bullets list evenly across the placeholders so layouts like
         "Comparison" (Before / After) actually get used as intended instead of
         dumping all bullets into the first column and overflowing.
      4. Otherwise distribute body_blocks 1:1 across text placeholders in order;
         if more blocks than placeholders, concatenate the overflow into the
         LAST placeholder (preserves content even when layout shape mismatches).
      5. Clear any trailing unused text placeholders so the rendered slide
         doesn't show PowerPoint's default "Click to add text" prompt.
      6. Record every empty picture placeholder as ANNOTATED so the user knows
         where an image asset would go.

    Fixes the silent-content-loss bug from the executive-briefing × Madison
    render: slide 7's diagram + prose used to drop the prose (no second text
    slot); slide 8's bullets used to overflow (all in first slot, ignoring the
    second). Both now bind across the layout's actual placeholder structure.
    """
    content_phs, body_phs, picture_placeholders = _classify_placeholders(slide)
    # Preferred binding order: content slots first, then body/header slots.
    text_placeholders = content_phs + body_phs

    if not body_blocks:
        for ph in text_placeholders:
            _safe_clear_text_frame(ph)
        return

    if not text_placeholders:
        manifest.add(
            "DROPPED", slide_id, "body_blocks",
            f"layout has no text-capable body placeholders; "
            f"{len(body_blocks)} body block(s) were dropped.",
        )
        return

    n_ph = len(text_placeholders)
    n_content = len(content_phs)
    n_blocks = len(body_blocks)

    # --- (3) Special case: single bullets block + 2+ CONTENT slots ---
    # Layouts like Comparison have two large OBJECT slots for left/right
    # columns. When IR has a single bullets block with enough items to fill
    # both columns, split the bullets across the content slots only — and
    # clear the small BODY/header slots so they don't show prompt text.
    if (n_blocks == 1 and n_content >= 2
            and body_blocks[0].get("kind") == "bullets"):
        bullets = body_blocks[0].get("content", [])
        if isinstance(bullets, list) and len(bullets) >= n_content:
            chunks = _split_evenly(bullets, n_content)
            for ph, chunk in zip(content_phs, chunks):
                _set_text_frame_to_bullets(ph, chunk)
            for ph in body_phs:
                _safe_clear_text_frame(ph)
            manifest.add(
                "ANNOTATED", slide_id, "body_blocks[0].kind=bullets",
                f"{len(bullets)} bullet(s) split across {n_content} "
                f"layout column(s).",
            )
            for ph in picture_placeholders:
                manifest.add(
                    "ANNOTATED", slide_id, "empty_picture_placeholder",
                    "picture placeholder left empty; IR supplied no image asset.",
                )
            return

    # --- (4) Standard distribute: one block per placeholder, concat overflow ---
    if n_blocks <= n_ph:
        # 1:1 binding; trailing placeholders cleared
        for i in range(n_blocks):
            block = body_blocks[i]
            text = _render_block_to_text(
                block.get("kind", ""), block.get("content"),
                manifest, slide_id, i,
            )
            _set_text_frame_to_text(text_placeholders[i], text)
        for ph in text_placeholders[n_blocks:]:
            _safe_clear_text_frame(ph)
    else:
        # First n_ph-1 placeholders get one block each; last gets the remainder
        for i in range(n_ph - 1):
            block = body_blocks[i]
            text = _render_block_to_text(
                block.get("kind", ""), block.get("content"),
                manifest, slide_id, i,
            )
            _set_text_frame_to_text(text_placeholders[i], text)
        # Concat the remaining blocks into the last placeholder
        last_ph = text_placeholders[-1]
        try:
            tf = last_ph.text_frame
        except Exception:
            tf = None
        if tf is not None:
            tf.clear()
            first_para = True
            for j in range(n_ph - 1, n_blocks):
                block = body_blocks[j]
                text = _render_block_to_text(
                    block.get("kind", ""), block.get("content"),
                    manifest, slide_id, j,
                )
                if text is None:
                    continue
                if first_para:
                    tf.paragraphs[0].text = text
                    first_para = False
                else:
                    p = tf.add_paragraph()
                    p.text = text

    # --- (6) Record empty picture placeholders for transparency ---
    for ph in picture_placeholders:
        manifest.add(
            "ANNOTATED", slide_id, "empty_picture_placeholder",
            "picture placeholder left empty; IR supplied no image asset.",
        )


# ---------------------------------------------------------------------------
# Placeholder classification + safe text helpers (v0.2.1)
# ---------------------------------------------------------------------------

def _classify_placeholders(slide: Any) -> tuple[list, list, list]:
    """Split non-title placeholders into three buckets:

      content_phs: OBJECT-typed (7) — the large content slots in layouts like
                   "Title and Content" and Comparison's two-column body areas.
                   Preferred for IR body_block binding.
      body_phs:    BODY-typed (2) — text placeholders. In single-body layouts
                   like "Title and Content" these ARE the body slot; in
                   multi-section layouts like Comparison they're typically
                   small header labels ("Before:"/"After:"). Used as overflow
                   or fallback when no OBJECT slots exist.
      picture_phs: media placeholders that can't hold text. Left empty;
                   recorded as ANNOTATED.

    System placeholders (DATE=16, FOOTER=15, SLIDE_NUMBER=13, HEADER=14) are
    excluded entirely — they hold metadata, not IR content.
    """
    SYSTEM_TYPES = {13, 14, 15, 16}
    OBJECT_TYPE = 7
    BODY_TYPE = 2

    content_phs, body_phs, picture_phs = [], [], []
    for ph in slide.placeholders:
        pf = ph.placeholder_format
        if pf.idx == 0:
            continue
        try:
            type_val = int(pf.type) if pf.type is not None else None
        except Exception:
            type_val = None
        if type_val in SYSTEM_TYPES:
            continue
        try:
            _ = ph.text_frame
        except Exception:
            picture_phs.append(ph)
            continue
        if type_val == OBJECT_TYPE:
            content_phs.append(ph)
        else:
            body_phs.append(ph)
    return content_phs, body_phs, picture_phs


def _safe_clear_text_frame(ph: Any) -> None:
    """Best-effort clear of a placeholder's text frame; ignore failures."""
    try:
        ph.text_frame.clear()
    except Exception:
        pass


def _set_text_frame_to_text(ph: Any, text: str | None) -> None:
    """Replace a placeholder's text frame contents with a single string."""
    if text is None:
        _safe_clear_text_frame(ph)
        return
    try:
        tf = ph.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
    except Exception:
        pass


def _set_text_frame_to_bullets(ph: Any, bullets: list[str]) -> None:
    """Replace a placeholder's text frame contents with a bullet list."""
    try:
        tf = ph.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = str(b)
            else:
                p = tf.add_paragraph()
                p.text = str(b)
    except Exception:
        pass


def _split_evenly(items: list, n_chunks: int) -> list[list]:
    """Distribute `items` across `n_chunks` lists as evenly as possible.

    For 5 items into 2 chunks → [3, 2]. For 7 items into 3 chunks → [3, 2, 2].
    Order is preserved within and across chunks.
    """
    if n_chunks <= 0:
        return []
    base, extra = divmod(len(items), n_chunks)
    out = []
    idx = 0
    for c in range(n_chunks):
        size = base + (1 if c < extra else 0)
        out.append(items[idx:idx + size])
        idx += size
    return out


# Body-block text rendering moved to _common.render_block_to_text (v0.2.2).
# Local alias preserved so existing call sites read cleanly.
from _common import render_block_to_text as _render_block_to_text  # noqa: E402


def set_speaker_notes(slide: Any, notes: str, manifest: LossManifest, slide_id: str) -> None:
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes
    manifest.add("LOSSLESS", slide_id, "speaker_notes",
                 f"speaker notes ({len(notes)} chars) preserved.")


def record_deck_metadata(ir: dict, manifest: LossManifest) -> None:
    """Record what was preserved/dropped at the deck level."""
    deck = ir.get("deck", {})
    # Deck-level fields that don't have a natural pptx home
    deck_fields_dropped = ["audience", "throughline", "arc", "evidence_anchor",
                            "duration_min", "voice_constraints"]
    for f in deck_fields_dropped:
        val = deck.get(f)
        if val:
            manifest.add(
                "DROPPED", None, f"deck.{f}",
                f"'{f}' has no native .pptx representation; preserved only in "
                f"the loss manifest (value: {str(val)[:80]}).",
            )
    if deck.get("title"):
        manifest.add("LOSSLESS", None, "deck.title",
                     "deck title applied to pptx core properties.")


# ---------------------------------------------------------------------------
# Main render
# ---------------------------------------------------------------------------

def render(ir_path: Path, profile_path: Path, out_path: Path) -> LossManifest:
    ir = load_ir(ir_path)
    profile = load_profile(profile_path)

    pptx_branch = profile.get("templates", {}).get("pptx")
    if not pptx_branch:
        raise ValueError("Profile has no `templates.pptx` branch. "
                         "render-pptx cannot proceed.")
    template_path_str = str(Path(pptx_branch["path"]).expanduser())
    layout_map = pptx_branch.get("layout_map", {})

    manifest = _make_manifest(
        deck_title=ir.get("deck", {}).get("title", ""),
        rendered_at=dt.datetime.now().isoformat(timespec="seconds"),
        template_path=template_path_str,
    )

    prs = Presentation(template_path_str)

    # Remove any existing slides from the template (we render into a fresh
    # deck inheriting only the template's masters/layouts/styles).
    _strip_existing_slides(prs)

    record_deck_metadata(ir, manifest)
    if prs.core_properties:
        prs.core_properties.title = ir.get("deck", {}).get("title", "")

    for slide_def in ir.get("slides", []):
        slide_id = slide_def.get("id", "")
        intent = slide_def.get("layout_intent", "")
        layout = resolve_layout(prs, layout_map, intent, manifest, slide_id)
        slide = prs.slides.add_slide(layout)

        set_title(slide, slide_def.get("title", ""), manifest, slide_id,
                  intent=slide_def.get("layout_intent", "claim_with_evidence"))
        populate_body(slide, slide_def.get("body_blocks", []), manifest, slide_id)
        set_speaker_notes(slide, slide_def.get("speaker_notes", ""), manifest, slide_id)

        # Transitions — IR-only concept, dropped in pptx
        transitions = slide_def.get("transitions")
        if transitions:
            manifest.add(
                "DROPPED", slide_id, "transitions",
                "transitions are an IR-only narrative-connective-tissue field; "
                "no native pptx representation. Visible in speaker_notes if "
                "the IR author copied them there.",
            )

    out_path = Path(out_path).expanduser()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    return manifest


def _strip_existing_slides(prs: Presentation) -> None:
    """Remove all slides from the prs, leaving masters/layouts intact.

    Two steps are required: (1) remove the slide ID from the presentation's
    sldIdLst, and (2) drop the relationship that connects the presentation
    part to the slide part. Without (2), the orphaned slide parts remain in
    the package and python-pptx writes them on save (triggering 'Duplicate
    name: ppt/slides/slideN.xml' warnings and bloating the file).
    """
    REL_ID_ATTR = (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    xml_slides = prs.slides._sldIdLst
    slides_to_remove = list(xml_slides)
    for slide_id in slides_to_remove:
        rId = slide_id.get(REL_ID_ATTR)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass  # already dropped
        xml_slides.remove(slide_id)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Render a Deck IR YAML into a .pptx using a template profile. "
            "Emits the .pptx plus a markdown loss manifest and JSON sidecar."
        )
    )
    parser.add_argument("--ir", required=True, help="Path to Deck IR YAML")
    parser.add_argument("--profile", required=True, help="Path to template profile YAML")
    parser.add_argument("--out", required=True, help="Output .pptx path")
    args = parser.parse_args()

    ir_path = Path(args.ir).expanduser()
    profile_path = Path(args.profile).expanduser()
    out_path = Path(args.out).expanduser()

    try:
        manifest = render(ir_path, profile_path, out_path)
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        return 1

    loss_md_path = out_path.with_suffix(out_path.suffix + ".loss.md")
    loss_json_path = out_path.with_suffix(out_path.suffix + ".loss.json")
    loss_md_path.write_text(manifest.to_markdown())
    loss_json_path.write_text(json.dumps(manifest.to_json(), indent=2) + "\n")

    print(f"Wrote {out_path}")
    print(f"Wrote {loss_md_path}")
    print(f"Wrote {loss_json_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
