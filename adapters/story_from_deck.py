"""
story_from_deck.py — reverse-engineer a Deck IR from an existing .pptx.

Stage 1 (slide-level extraction): walks every slide, captures title, body
text, speaker notes, layout name, basic geometry.

Stage 2 (beat inference): tags each slide with a candidate beat name based
on position + layout-intent mapping.

Stage 3-4 (throughline + body_blocks): heuristics over text.

Stage 5 (style influence): adjusts arc + beat allocation per the chosen
style file.

Stage 6 (emit IR + caveats sidecar).

v0.2 scope. Style-aware. Companion to `story-from-deck` SKILL.md.

Anonymity: ships in the public plugin. No organization patterns; no
hard-coded user content.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from collections import Counter
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).parent))
from _common import (  # noqa: E402
    LAYOUT_INTENTS, BEATS, BODY_BLOCK_KINDS, ARC_VALUES, EVIDENCE_ANCHORS,
    detect_intent_from_name,
)

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx required. pip install python-pptx", file=sys.stderr)
    sys.exit(2)

try:
    import yaml
except ImportError:
    print("ERROR: pyyaml required. pip install pyyaml", file=sys.stderr)
    sys.exit(2)


# ---- Beat inference --------------------------------------------------------

# Position-based beat priors. First slide is almost always `hook`;
# last is almost always `next` or `callback`. Middle slides default
# to `evidence` until layout-intent overrides.
def _position_beat(slide_idx: int, n_slides: int) -> str:
    if slide_idx == 0:
        return "hook"
    if slide_idx == 1 and n_slides > 5:
        return "context"
    if slide_idx == n_slides - 1:
        return "next"
    if slide_idx == n_slides - 2:
        return "callback"
    # Middle: distribute across claim / evidence / tension / resolution
    mid_third = (slide_idx - 1) / max(1, n_slides - 2)
    if mid_third < 0.25:
        return "context"
    if mid_third < 0.5:
        return "claim"
    if mid_third < 0.85:
        return "evidence"
    return "resolution"


# Reverse map: intent → preferred beat (most common case)
_INTENT_TO_LIKELY_BEAT: dict[str, str] = {
    "title": "hook",
    "section_break": "context",
    "claim_with_evidence": "claim",
    "three_pillars": "resolution",
    "comparison": "tension",
    "quote": "evidence",
    "image_with_caption": "evidence",
    "metrics": "evidence",
    "timeline": "context",
    "callout": "callback",
}


def infer_beat(slide_idx: int, n_slides: int, layout_name: str) -> str:
    """Combine position prior + layout-intent reverse lookup."""
    intent = detect_intent_from_name(layout_name)
    if intent and intent in _INTENT_TO_LIKELY_BEAT:
        return _INTENT_TO_LIKELY_BEAT[intent]
    return _position_beat(slide_idx, n_slides)


# ---- Slide extraction ------------------------------------------------------

def extract_slide(slide, layout_name: str, idx: int) -> dict[str, Any]:
    """Capture title, body text, speaker notes, beat candidate."""
    title = ""
    if slide.shapes.title is not None:
        try:
            title = slide.shapes.title.text or ""
        except Exception:
            title = ""

    body_text_parts: list[str] = []
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            continue  # already captured as title
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if text:
            body_text_parts.append(text)
    # Also pick up text from non-placeholder shapes (text boxes etc.)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            # is_placeholder raises AttributeError on some shape types;
            # guard by checking the property directly.
            if getattr(shape, "is_placeholder", False):
                continue
        except Exception:
            continue
        text = (shape.text_frame.text or "").strip()
        if text and text != title and text not in body_text_parts:
            body_text_parts.append(text)
    body_text = "\n\n".join(body_text_parts)

    speaker_notes = ""
    try:
        if slide.has_notes_slide:
            speaker_notes = (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        pass

    return {
        "slide_idx": idx,
        "title": title.strip(),
        "body_text": body_text,
        "speaker_notes": speaker_notes,
        "layout_name": layout_name,
    }


# ---- Body block reconstruction --------------------------------------------

_METRIC_RE = re.compile(r"(\$?\d[\d,.]*\s*[%a-zA-Z]+)")  # rough: "47%", "$8.2M", "12 days"


def _looks_like_bullets(text: str) -> bool:
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    if len(lines) < 2:
        return False
    bullet_prefixes = ("•", "-", "*", "·", "—", "–")
    bulleted = sum(1 for l in lines if l.startswith(bullet_prefixes))
    return bulleted >= len(lines) // 2


def _looks_like_quote(text: str) -> bool:
    return (text.startswith(('"', '"', "'", "“")) and
            text.endswith(('"', '"', "'", "”")))


def reconstruct_body_blocks(body_text: str) -> list[dict[str, Any]]:
    """Convert raw body text into IR body_blocks."""
    if not body_text or not body_text.strip():
        return []
    blocks: list[dict[str, Any]] = []
    if _looks_like_bullets(body_text):
        lines = [
            re.sub(r"^[\s•\-\*·—–]+", "", l).strip()
            for l in body_text.split("\n") if l.strip()
        ]
        blocks.append({"kind": "bullets", "content": lines})
        return blocks
    # Check for metric pattern in short text
    if len(body_text) < 200:
        metric_match = _METRIC_RE.search(body_text)
        if metric_match and len(body_text.split()) <= 30:
            value = metric_match.group(1)
            label = body_text.replace(value, "").strip(" -—:")
            blocks.append({
                "kind": "metric",
                "content": {"label": label or "Result", "value": value},
            })
            return blocks
    # Quote detection
    if _looks_like_quote(body_text):
        text = body_text.strip(' "\'""''')
        blocks.append({"kind": "quote", "content": {"text": text}})
        return blocks
    # Default: prose
    blocks.append({"kind": "prose", "content": body_text})
    return blocks


# ---- Throughline inference -------------------------------------------------

def infer_throughline(slides: list[dict[str, Any]]) -> str:
    """Most-repeated meaningful 4-7 word phrase across titles + notes.

    v0.1 heuristic: look at slide titles, pick the longest non-trivial one
    that resembles a claim. If no good candidate, return empty (the skill's
    caller will prompt the user).
    """
    candidates: list[str] = []
    for s in slides:
        title = s.get("title", "")
        if 6 < len(title) < 150 and not title.lower().startswith(("part ", "chapter ")):
            candidates.append(title)
    # Prefer titles with action verbs / claims (rough heuristic)
    claim_keywords = ("will", "must", "should", "is", "becomes", "wins", "drives", "makes")
    for c in candidates:
        if any(kw in c.lower().split() for kw in claim_keywords):
            return c
    if candidates:
        return candidates[len(candidates) // 2]  # middle of the deck
    return ""


# ---- IR assembly -----------------------------------------------------------

def build_ir(slides_data: list[dict[str, Any]],
             style: dict[str, Any] | None) -> tuple[dict[str, Any], list[str]]:
    n = len(slides_data)
    caveats: list[str] = []

    throughline = infer_throughline(slides_data)
    if not throughline:
        caveats.append("Could not infer a throughline from slide content. "
                       "User must provide one before locking the IR.")

    # Arc: default from style, or 'situation-first'
    arc = "situation-first"
    if style:
        arc = style.get("default_arc", arc)

    # Evidence anchor: count body-block kinds and pick the dominant
    kinds_seen: Counter[str] = Counter()
    for s in slides_data:
        for block in reconstruct_body_blocks(s.get("body_text", "")):
            kinds_seen[block["kind"]] += 1
    if kinds_seen.get("metric", 0) >= 2:
        evidence_anchor = "numbers"
    elif kinds_seen.get("quote", 0) >= 2:
        evidence_anchor = "story"
    elif kinds_seen.get("image_placeholder", 0) >= 1:
        evidence_anchor = "demo"
    else:
        evidence_anchor = "hybrid"

    ir_slides: list[dict[str, Any]] = []
    for s in slides_data:
        idx = s["slide_idx"]
        beat = infer_beat(idx, n, s.get("layout_name", ""))
        intent = detect_intent_from_name(s.get("layout_name", "")) or "claim_with_evidence"
        # Title fallback chain: title placeholder → first line of body → "Slide N"
        title = s.get("title") or ""
        if not title.strip():
            first_body_line = (s.get("body_text") or "").split("\n", 1)[0].strip()
            title = first_body_line[:80] if first_body_line else f"Slide {idx + 1}"
        ir_slide: dict[str, Any] = {
            "id": f"slide-{idx+1:02d}" + (f"-{_slug(title)}" if title else ""),
            "beat": beat,
            "layout_intent": intent,
            "title": title,
        }
        blocks = reconstruct_body_blocks(s.get("body_text", ""))
        if blocks:
            ir_slide["body_blocks"] = blocks
        if s.get("speaker_notes"):
            ir_slide["speaker_notes"] = s["speaker_notes"]
        ir_slides.append(ir_slide)

    deck_title = slides_data[0]["title"] if slides_data and slides_data[0].get("title") else "Recovered deck"
    ir = {
        "ir_version": "1.0.0",
        "deck": {
            "title": deck_title,
            "audience": {
                "primary": "(user-supplied)",
                "prior_knowledge": "warm",
            },
            "throughline": throughline or "(user-supplied)",
            "arc": arc,
            "evidence_anchor": evidence_anchor,
        },
        "slides": ir_slides,
    }
    if throughline == "":
        caveats.append("deck.throughline left as '(user-supplied)' — fill in before render.")
    caveats.append("deck.audience.primary left as '(user-supplied)' — fill in before render.")

    return ir, caveats


def _slug(text: str) -> str:
    """Kebab-case slug from text, max 40 chars."""
    s = re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")
    return s[:40]


# ---- Main -----------------------------------------------------------------

def reverse_engineer(pptx_path: Path, style_path: Path | None) -> tuple[dict[str, Any], list[str]]:
    prs = Presentation(str(pptx_path))
    slides_data: list[dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else ""
        slides_data.append(extract_slide(slide, layout_name, i))

    style = None
    if style_path and style_path.exists():
        style = yaml.safe_load(style_path.read_text())
    return build_ir(slides_data, style)


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Reverse-engineer a Deck IR from an existing .pptx."
    )
    parser.add_argument("pptx", help="Path to source .pptx")
    parser.add_argument("--style", help="Path to a style YAML (e.g., styles/ted-talk.yaml)")
    parser.add_argument("--out", required=True, help="Output IR YAML path")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).expanduser()
    if not pptx_path.exists():
        print(f"ERROR: {pptx_path} not found.", file=sys.stderr)
        return 2

    style_path = Path(args.style).expanduser() if args.style else None
    try:
        ir, caveats = reverse_engineer(pptx_path, style_path)
    except Exception as e:
        print(f"ERROR: {e}", file=sys.stderr)
        return 1

    out_path = Path(args.out).expanduser()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(yaml.safe_dump(ir, sort_keys=False, allow_unicode=True))

    # Sidecar caveats
    if caveats:
        sidecar = out_path.with_suffix(out_path.suffix + ".reverse-engineering.md")
        sidecar.write_text(
            "# Reverse-engineering caveats\n\n"
            f"Recovered from: `{pptx_path}`\n\n"
            "## Items to review\n\n"
            + "\n".join(f"- {c}" for c in caveats) + "\n"
        )
        print(f"Wrote {sidecar}")

    print(f"Wrote {out_path}")
    print(f"Slides recovered: {len(ir['slides'])}")
    print(f"Throughline: {ir['deck']['throughline']}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
