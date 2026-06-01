"""
template_synthesizer_pptx.py — synthesize a candidate template from a
deck-with-implicit-pattern .pptx.

Inputs:
    deck_path:  path to the source .pptx (classified as
                "deck-with-implicit-pattern" by template-classifier)
    out_path:   where to write the synthesized .pptx
    report_path: where to write the synthesizer-report JSON

Outputs:
    <out_path>.pptx     — synthesized template (renamed layouts; v0.1 v0.1
                         carries layout NAMES + token info, not canonical
                         shape positions; v0.2 will embed shapes)
    <report_path>.json  — derived patterns + suggested layout_map + tokens

Pipeline:
    1. Profile each slide: shape count, geometry signature, font, color.
    2. Cluster slides by signature (exact equality on quantized signature).
    3. For each cluster, derive a canonical pattern + name via heuristics.
    4. Build a synthesized .pptx by renaming default layouts to match clusters.
    5. Emit a JSON report with the per-cluster patterns + tokens for the user
       to manually rebuild a proper template (or for v0.2 shape embedding).

Anonymity: ships in the public plugin. Heuristics are generic.
"""

from __future__ import annotations

import argparse
import json
import sys
from collections import Counter, defaultdict
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
except ImportError:
    print("ERROR: python-pptx is not installed. Install with: "
          "pip install python-pptx", file=sys.stderr)
    sys.exit(2)


# Position quantization granularity (in EMUs — 914400 EMU = 1 inch).
# Half-inch grid is coarse enough to forgive small position drift.
POSITION_QUANTUM = 457200  # 0.5 inch in EMUs

# Font-size quantization (pt) — group nearby sizes together.
FONT_SIZE_QUANTUM = 4


def _quantize(val: int | float | None, q: int) -> int:
    if val is None:
        return 0
    return int(round(val / q) * q)


@dataclass
class ShapeProfile:
    """A single shape's signature for cluster comparison."""
    kind: str           # "text" | "image" | "other"
    left: int           # quantized EMU
    top: int            # quantized EMU
    width: int          # quantized EMU
    height: int         # quantized EMU
    dominant_font_pt: int    # quantized
    dominant_font_family: str
    dominant_fill: str  # "#RRGGBB" or "" if none

    def signature_key(self) -> tuple:
        return (self.kind, self.left, self.top, self.width, self.height,
                self.dominant_font_pt)


@dataclass
class SlideProfile:
    """All shapes on one slide, ordered by position."""
    slide_idx: int
    shapes: list[ShapeProfile] = field(default_factory=list)

    def signature_key(self) -> tuple:
        return tuple(s.signature_key() for s in self.shapes)


@dataclass
class Cluster:
    """A group of slides sharing a signature."""
    cluster_id: int
    derived_name: str
    member_slide_indices: list[int]
    canonical_shapes: list[ShapeProfile]
    suggested_intent: str   # one of IR's 10 intents, or "unknown"


def _shape_kind(shape: Any) -> str:
    if shape.has_text_frame:
        return "text"
    try:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            return "image"
    except Exception:
        pass
    return "other"


def _dominant_font(shape: Any) -> tuple[int, str]:
    """Return (size_pt, family) for the largest font run in the shape."""
    if not shape.has_text_frame:
        return (0, "")
    sizes = []
    families = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
            if run.font.name:
                families.append(run.font.name)
    if not sizes and not families:
        return (0, "")
    size = max(sizes) if sizes else 0
    family = Counter(families).most_common(1)[0][0] if families else ""
    return (_quantize(int(size), FONT_SIZE_QUANTUM), family)


def _dominant_fill(shape: Any) -> str:
    """Return the dominant text run color for text shapes; '' otherwise."""
    if not shape.has_text_frame:
        return ""
    colors = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.color and run.font.color.rgb is not None:
                    colors.append(str(run.font.color.rgb))
            except Exception:
                continue
    if not colors:
        return ""
    return "#" + Counter(colors).most_common(1)[0][0].upper()


def profile_slide(slide: Any, idx: int) -> SlideProfile:
    shapes: list[ShapeProfile] = []
    for sh in slide.shapes:
        try:
            left = _quantize(sh.left, POSITION_QUANTUM) if sh.left is not None else 0
            top = _quantize(sh.top, POSITION_QUANTUM) if sh.top is not None else 0
            width = _quantize(sh.width, POSITION_QUANTUM) if sh.width is not None else 0
            height = _quantize(sh.height, POSITION_QUANTUM) if sh.height is not None else 0
        except Exception:
            continue
        font_pt, font_family = _dominant_font(sh)
        fill = _dominant_fill(sh)
        shapes.append(ShapeProfile(
            kind=_shape_kind(sh),
            left=left, top=top, width=width, height=height,
            dominant_font_pt=font_pt,
            dominant_font_family=font_family,
            dominant_fill=fill,
        ))
    # Sort shapes by (top, left) for stable comparison
    shapes.sort(key=lambda s: (s.top, s.left))
    return SlideProfile(slide_idx=idx, shapes=shapes)


def cluster_slides(profiles: list[SlideProfile]) -> list[Cluster]:
    """Group slides by exact signature_key. First-seen ordering."""
    groups: dict[tuple, list[int]] = defaultdict(list)
    canonical: dict[tuple, list[ShapeProfile]] = {}
    for p in profiles:
        key = p.signature_key()
        groups[key].append(p.slide_idx)
        if key not in canonical:
            canonical[key] = p.shapes
    clusters: list[Cluster] = []
    for cid, (key, members) in enumerate(groups.items()):
        shapes = canonical[key]
        name = _derive_layout_name(shapes)
        intent = _suggest_intent(name, shapes)
        clusters.append(Cluster(
            cluster_id=cid,
            derived_name=name,
            member_slide_indices=members,
            canonical_shapes=shapes,
            suggested_intent=intent,
        ))
    return clusters


def _derive_layout_name(shapes: list[ShapeProfile]) -> str:
    """Heuristic naming based on shape count, position, and font sizes."""
    n_text = sum(1 for s in shapes if s.kind == "text")
    n_image = sum(1 for s in shapes if s.kind == "image")
    if not shapes:
        return "Blank Layout"

    # Inspect text shape font sizes (largest first)
    text_sizes = sorted([s.dominant_font_pt for s in shapes if s.kind == "text"],
                       reverse=True)
    largest = text_sizes[0] if text_sizes else 0

    # 1 big text + 1 small. Order matters — check huge-font cases first.
    if n_text == 2 and n_image == 0:
        if largest >= 80:
            # Huge-font slides are metric/big-statement layouts, not titles.
            return "Stat Block"
        if largest >= 40 and text_sizes[1] <= 24:
            return "Title Slide"
        if largest >= 24 and text_sizes[1] <= 16:
            return "Pull Quote"

    # Header + 3 columns of text → Three Column
    if n_text == 4 and n_image == 0:
        # Top shape is the header, 3 below are columns
        cols = shapes[-3:]
        # If the 3 columns have similar widths, classify as Three Column
        widths = sorted({s.width for s in cols})
        if len(widths) <= 2:
            return "Three Column"

    # 1 huge number + 1 small label → Stat Block / Metric
    if n_text == 2 and largest >= 80:
        return "Stat Block"

    # Single text shape → Section Header
    if n_text == 1 and n_image == 0:
        return "Section Header"

    # Image-heavy
    if n_image >= 1 and n_text <= 2:
        return "Image and Caption"

    return f"Custom Layout (n_text={n_text}, n_image={n_image})"


def _suggest_intent(layout_name: str, shapes: list[ShapeProfile]) -> str:
    """Map the derived layout name to an IR intent."""
    name = layout_name.lower()
    if "title slide" in name:
        return "title"
    if "section" in name:
        return "section_break"
    if "three" in name or "column" in name:
        return "three_pillars"
    if "quote" in name:
        return "quote"
    if "stat" in name or "metric" in name or "big statement" in name:
        return "metrics" if "stat" in name or "metric" in name else "callout"
    if "image" in name or "picture" in name:
        return "image_with_caption"
    if "compar" in name:
        return "comparison"
    if "timeline" in name:
        return "timeline"
    return "claim_with_evidence"  # universal default


def extract_global_tokens(profiles: list[SlideProfile]) -> dict[str, Any]:
    """Aggregate token-like signals across all profiled slides."""
    families: Counter[str] = Counter()
    sizes: Counter[int] = Counter()
    fills: Counter[str] = Counter()
    for p in profiles:
        for s in p.shapes:
            if s.dominant_font_family:
                families[s.dominant_font_family] += 1
            if s.dominant_font_pt > 0:
                sizes[s.dominant_font_pt] += 1
            if s.dominant_fill:
                fills[s.dominant_fill] += 1

    color_roles = ["primary", "secondary", "accent", "text-primary",
                   "text-secondary", "surface", "surface-muted"]
    colors_out: dict[str, str] = {}
    for role, (hex_val, _) in zip(color_roles, fills.most_common(len(color_roles))):
        colors_out[role] = hex_val

    primary_family = families.most_common(1)[0][0] if families else ""
    type_roles = ["display", "heading-1", "heading-2", "body", "caption"]
    distinct_sizes = sorted({s for s, _ in sizes.most_common(20) if s > 0},
                            reverse=True)
    typography_out: dict[str, dict[str, Any]] = {}
    for i, role in enumerate(type_roles):
        if i >= len(distinct_sizes):
            break
        typography_out[role] = {
            "family": primary_family or "Helvetica",
            "size_pt": float(distinct_sizes[i]),
            "weight": 700 if i < 2 else 400,
        }

    return {"colors": colors_out, "typography": typography_out}


def build_synthesized_pptx(clusters: list[Cluster], out_path: Path) -> None:
    """Build a synthesized .pptx whose layouts are renamed to match clusters.

    v0.1: this gives the renderer a template with the right LAYOUT NAMES.
    Canonical shape positions are documented in the JSON report; v0.2 will
    embed them into the layouts.
    """
    prs = Presentation()
    master = prs.slide_masters[0]
    available_layouts = list(master.slide_layouts)
    for i, cluster in enumerate(clusters):
        if i >= len(available_layouts):
            break
        available_layouts[i].element.cSld.set("name", cluster.derived_name)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def synthesize(deck_path: Path, out_pptx_path: Path,
               report_path: Path | None = None) -> dict[str, Any]:
    prs = Presentation(str(deck_path))
    profiles = [profile_slide(s, i) for i, s in enumerate(prs.slides)]
    clusters = cluster_slides(profiles)
    tokens = extract_global_tokens(profiles)

    # Build suggested layout_map (IR intent → derived layout name)
    layout_map: dict[str, str] = {}
    for c in clusters:
        intent = c.suggested_intent
        # First cluster per intent wins
        if intent not in layout_map:
            layout_map[intent] = c.derived_name

    build_synthesized_pptx(clusters, out_pptx_path)

    report = {
        "source_deck": str(deck_path),
        "synthesized_template": str(out_pptx_path),
        "n_slides": len(profiles),
        "n_clusters": len(clusters),
        "clusters": [
            {
                "cluster_id": c.cluster_id,
                "derived_name": c.derived_name,
                "suggested_intent": c.suggested_intent,
                "member_slide_indices": c.member_slide_indices,
                "n_members": len(c.member_slide_indices),
                "canonical_shapes": [asdict(s) for s in c.canonical_shapes],
            }
            for c in clusters
        ],
        "suggested_layout_map": layout_map,
        "extracted_tokens": tokens,
        "v0_1_caveats": [
            "Synthesized .pptx carries layout NAMES only; canonical shape "
            "positions are documented in this report but not embedded in "
            "the .pptx layouts. v0.2 will embed shapes into the layouts.",
            "Renderer using this synthesized template will produce decks "
            "with the right structural skeleton but without the source "
            "deck's visual styling. The user can either (a) accept this "
            "as a barebones template and add styling in PowerPoint, or "
            "(b) use the canonical shape data in this report to manually "
            "rebuild the template with correct positions/fonts/colors.",
        ],
    }
    if report_path:
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, indent=2) + "\n")
    return report


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Synthesize a candidate template from a deck-with-implicit-pattern "
            ".pptx. Emits a synthesized .pptx (renamed layouts) and a JSON "
            "report (cluster patterns + tokens for v0.2 / user rebuild)."
        )
    )
    parser.add_argument("deck", help="Path to source .pptx deck")
    parser.add_argument("--out", required=True, help="Output .pptx path")
    parser.add_argument("--report", help="Output JSON report path "
                        "(default: <out>.synthesizer-report.json)")
    args = parser.parse_args()

    deck_path = Path(args.deck).expanduser()
    if not deck_path.exists():
        print(f"ERROR: {deck_path} not found.", file=sys.stderr)
        return 2

    out_path = Path(args.out).expanduser()
    report_path = (Path(args.report).expanduser() if args.report
                   else out_path.with_suffix(".synthesizer-report.json"))

    try:
        report = synthesize(deck_path, out_path, report_path)
    except Exception as e:
        print(f"ERROR synthesizing: {e}", file=sys.stderr)
        return 1

    print(f"Wrote {out_path}")
    print(f"Wrote {report_path}")
    print(f"Clusters: {report['n_clusters']} from {report['n_slides']} slides")
    for c in report["clusters"]:
        print(f"  - cluster {c['cluster_id']}: '{c['derived_name']}' "
              f"(intent: {c['suggested_intent']}, members: {c['n_members']})")
    return 0


if __name__ == "__main__":
    sys.exit(main())
