"""
template_classifier.py — detect whether a .pptx file is a structured template
or a deck-with-implicit-pattern.

Inputs:
    pptx_path: path to a .pptx file

Outputs:
    {
        "format": "pptx",
        "classification": "template" | "deck-with-implicit-pattern" | "mixed",
        "confidence": 0.0-1.0,
        "signals": {
            "n_slides": int,
            "n_layouts_available": int,
            "n_layouts_used": int,
            "layout_diversity_ratio": float,  # n_layouts_used / n_slides
            "default_layout_ratio": float,    # fraction on Title/Title-and-Content/Blank
            "direct_overrides_estimate": int, # shapes added to slides past placeholder count
            "layout_name_semantic_richness": float,  # 0-1, how many layout names match IR intent patterns
        },
        "diagnosis": [<human-readable observation>, ...],
    }

For figma, the analogous classifier runs as a Plugin API recipe (described in
the SKILL.md) and produces an equivalent JSON; this Python adapter handles the
pptx side only.

Anonymity: ships in the public plugin. No organization patterns, no defaults
baked in.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx is not installed. Install with: "
          "pip install python-pptx", file=sys.stderr)
    sys.exit(2)


# Layouts that signal "default / unmodified" usage. If most slides reference
# only these, the deck is using stock layouts rather than custom ones.
DEFAULT_LAYOUT_NAMES = {
    "title slide", "title", "title and content", "blank",
    "section header", "two content", "comparison",
    "title only", "content with caption", "picture with caption",
    "vertical title and text",
}

# IR layout intent patterns — same as the extractor uses. Match against
# layout names to estimate semantic richness.
INTENT_PATTERNS = [
    "title slide", "section", "claim", "three", "pillar",
    "comparison", "compare", "quote", "image", "picture",
    "stat", "metric", "kpi", "timeline", "callout", "big statement",
]


def classify_pptx(pptx_path: str | Path) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))

    # Collect layout metadata
    all_layouts: list[Any] = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            all_layouts.append(layout)

    layout_names = [l.name for l in all_layouts]
    n_layouts_available = len(all_layouts)

    # Map each slide to its layout
    slides = list(prs.slides)
    n_slides = len(slides)

    if n_slides == 0:
        return _empty_result(n_layouts_available, layout_names)

    used_layout_names: list[str] = []
    direct_overrides_total = 0
    default_layout_count = 0
    for slide in slides:
        layout = slide.slide_layout
        if layout is None:
            used_layout_names.append("(no layout)")
            continue
        used_layout_names.append(layout.name)
        if layout.name.lower() in DEFAULT_LAYOUT_NAMES:
            default_layout_count += 1

        # Estimate "direct overrides": shapes on the slide beyond what the
        # layout's placeholders provide. A hand-positioned text box adds to
        # this count; using a placeholder does not.
        layout_placeholder_count = sum(1 for ph in layout.placeholders)
        slide_shape_count = sum(1 for s in slide.shapes)
        if slide_shape_count > layout_placeholder_count:
            direct_overrides_total += slide_shape_count - layout_placeholder_count

    # Signals
    unique_used = set(used_layout_names)
    n_layouts_used = len(unique_used)
    layout_diversity_ratio = n_layouts_used / n_slides if n_slides else 0.0
    default_layout_ratio = default_layout_count / n_slides if n_slides else 0.0

    semantic_matches = sum(
        1 for name in layout_names
        if any(p in name.lower() for p in INTENT_PATTERNS)
    )
    layout_name_semantic_richness = (
        semantic_matches / n_layouts_available if n_layouts_available else 0.0
    )

    direct_overrides_per_slide = (
        direct_overrides_total / n_slides if n_slides else 0.0
    )

    signals = {
        "n_slides": n_slides,
        "n_layouts_available": n_layouts_available,
        "n_layouts_used": n_layouts_used,
        "layout_diversity_ratio": round(layout_diversity_ratio, 3),
        "default_layout_ratio": round(default_layout_ratio, 3),
        "direct_overrides_estimate": direct_overrides_total,
        "direct_overrides_per_slide": round(direct_overrides_per_slide, 2),
        "layout_name_semantic_richness": round(layout_name_semantic_richness, 3),
    }

    classification, confidence, diagnosis = _score(signals, layout_names, used_layout_names)

    return {
        "format": "pptx",
        "classification": classification,
        "confidence": confidence,
        "signals": signals,
        "diagnosis": diagnosis,
    }


def _empty_result(n_layouts: int, layout_names: list[str]) -> dict[str, Any]:
    return {
        "format": "pptx",
        "classification": "mixed",
        "confidence": 0.0,
        "signals": {
            "n_slides": 0, "n_layouts_available": n_layouts,
            "n_layouts_used": 0, "layout_diversity_ratio": 0.0,
            "default_layout_ratio": 0.0,
            "direct_overrides_estimate": 0,
            "direct_overrides_per_slide": 0.0,
            "layout_name_semantic_richness": 0.0,
        },
        "diagnosis": [
            "File has no slides; cannot classify. Likely an empty template "
            "shell — treat as 'template' for downstream extractor or have the "
            "user supply a populated file."
        ],
    }


def _score(signals: dict[str, Any], layout_names: list[str],
           used_layout_names: list[str]) -> tuple[str, float, list[str]]:
    """Apply heuristic rules to derive a classification + confidence."""
    diagnosis: list[str] = []

    # Key signals
    div = signals["layout_diversity_ratio"]
    default_ratio = signals["default_layout_ratio"]
    direct_overrides = signals["direct_overrides_per_slide"]
    semantic = signals["layout_name_semantic_richness"]
    n_used = signals["n_layouts_used"]

    # Score components (each contributes to a "template-ness" score)
    # template-ness in [0, 1]; higher = more template-like
    component_scores: dict[str, float] = {}

    # Diversity: high diversity → template (slides use varied layouts)
    component_scores["diversity"] = min(1.0, div * 2.0)  # 0.5 div → 1.0
    if div < 0.2:
        diagnosis.append(
            f"Layout diversity ratio is {div:.2f} — most slides share the "
            "same layout, suggesting a deck written on default layouts."
        )
    elif div > 0.4:
        diagnosis.append(
            f"Layout diversity ratio is {div:.2f} — slides use varied layouts, "
            "consistent with template-driven authoring."
        )

    # Default-layout ratio: high → deck (most slides on stock layouts)
    component_scores["non_default"] = 1.0 - default_ratio
    if default_ratio > 0.7:
        diagnosis.append(
            f"{int(default_ratio * 100)}% of slides reference stock layouts "
            "(Title Slide, Title and Content, Blank). Common in decks written "
            "on the default template without modification."
        )

    # Direct overrides: high → deck (hand-positioned shapes everywhere)
    if direct_overrides > 3:
        component_scores["low_overrides"] = max(0.0, 1.0 - (direct_overrides - 3) / 5)
        diagnosis.append(
            f"Estimated {direct_overrides:.1f} hand-positioned shapes per "
            "slide. High direct-override counts indicate content authored "
            "as text boxes / shapes rather than via placeholders."
        )
    else:
        component_scores["low_overrides"] = 1.0

    # Semantic richness of layout names: high → template
    component_scores["semantic_richness"] = semantic
    if semantic > 0.4:
        diagnosis.append(
            f"{int(semantic * 100)}% of layouts have semantically rich names "
            "(matching IR intents like 'Three Column', 'Pull Quote'). "
            "Consistent with intentional template design."
        )
    elif semantic < 0.2:
        diagnosis.append(
            f"Only {int(semantic * 100)}% of layouts have semantically "
            "meaningful names — most are generic ('Title and Content', "
            "'Blank', or unnamed). Templates designed for slide-publisher "
            "typically expose intent through layout names."
        )

    # n_layouts_used is a tie-breaker
    if n_used < 3:
        component_scores["layout_breadth"] = 0.0
    elif n_used >= 5:
        component_scores["layout_breadth"] = 1.0
    else:
        component_scores["layout_breadth"] = (n_used - 3) / 2

    # Weighted average
    weights = {
        "diversity": 0.30,
        "non_default": 0.25,
        "low_overrides": 0.20,
        "semantic_richness": 0.15,
        "layout_breadth": 0.10,
    }
    template_ness = sum(component_scores[k] * weights[k] for k in weights)

    # Confidence is distance from the middle (0.5)
    confidence = round(abs(template_ness - 0.5) * 2, 2)

    if template_ness >= 0.65:
        classification = "template"
        diagnosis.append(
            f"Verdict: 'template' (template-ness score {template_ness:.2f}, "
            f"confidence {confidence})."
        )
    elif template_ness <= 0.35:
        classification = "deck-with-implicit-pattern"
        diagnosis.append(
            f"Verdict: 'deck-with-implicit-pattern' (template-ness score "
            f"{template_ness:.2f}, confidence {confidence}). Recommend "
            "running template-synthesizer-pptx to derive a template from the "
            "recurring visual patterns in the deck."
        )
    else:
        classification = "mixed"
        diagnosis.append(
            f"Verdict: 'mixed' (template-ness score {template_ness:.2f}, "
            f"confidence {confidence}). Some template structure present but "
            "uneven. Manual review recommended; the extractor may produce a "
            "partial profile that the user should augment."
        )

    return classification, confidence, diagnosis


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Classify a .pptx as template, deck-with-implicit-pattern, or "
            "mixed. Emits JSON with the verdict, signals, and diagnosis."
        )
    )
    parser.add_argument("pptx", help="Path to .pptx")
    parser.add_argument("--out", help="Output JSON path (default: stdout)")
    args = parser.parse_args()

    p = Path(args.pptx).expanduser()
    if not p.exists():
        print(f"ERROR: {p} not found.", file=sys.stderr)
        return 2

    try:
        result = classify_pptx(p)
    except Exception as e:
        print(f"ERROR classifying: {e}", file=sys.stderr)
        return 1

    out_text = json.dumps(result, indent=2)
    if args.out:
        Path(args.out).expanduser().write_text(out_text + "\n")
    else:
        print(out_text)
    return 0


if __name__ == "__main__":
    sys.exit(main())
